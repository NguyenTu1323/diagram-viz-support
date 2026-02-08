#!/usr/bin/env python3
"""
Example PowerPoint Generator
Creates a data-driven presentation with charts and tables.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import pandas as pd
import io
from pathlib import Path

# Sample data
revenue_data = {
    'Month': ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
              'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec'],
    'Revenue': [45, 52, 48, 61, 58, 67, 72, 69, 75, 81, 78, 89]
}

performance_data = {
    'Metric': ['API Latency (ms)', 'Uptime (%)', 'Cache Hit Rate (%)', 'Error Rate (%)'],
    'Target': ['< 100', '99.9', '> 80', '< 0.1'],
    'Current': ['78', '99.95', '87', '0.04'],
    'Status': ['✅ Excellent', '✅ Exceeding', '✅ Optimal', '✅ Great']
}


def create_title_slide(prs):
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "2024 Platform Performance Report"
    subtitle.text = "Revenue Growth & System Metrics\nGenerated with Python"

    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)


def create_chart_slide(prs):
    """Create a slide with a bar chart using matplotlib."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Monthly Revenue Growth 2024"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    # Create matplotlib chart
    fig, ax = plt.subplots(figsize=(10, 5))
    df = pd.DataFrame(revenue_data)

    colors = ['#4e79a7', '#4e79a7', '#4e79a7',  # Q1
              '#59a14f', '#59a14f', '#59a14f',  # Q2
              '#f28e2c', '#f28e2c', '#f28e2c',  # Q3
              '#e15759', '#e15759', '#e15759']  # Q4

    bars = ax.bar(df['Month'], df['Revenue'], color=colors, edgecolor='white', linewidth=1.2)
    ax.set_xlabel('Month', fontsize=12, fontweight='bold')
    ax.set_ylabel('Revenue ($K)', fontsize=12, fontweight='bold')
    ax.set_ylim(0, 100)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    # Add value labels on bars
    for bar in bars:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'${int(height)}K',
                ha='center', va='bottom', fontsize=9, fontweight='bold')

    plt.tight_layout()

    # Save to BytesIO and add to slide
    image_stream = io.BytesIO()
    plt.savefig(image_stream, format='png', dpi=150, bbox_inches='tight')
    image_stream.seek(0)
    plt.close()

    slide.shapes.add_picture(image_stream, Inches(0.5), Inches(1.2),
                            width=Inches(9), height=Inches(4.5))


def create_table_slide(prs):
    """Create a slide with a performance metrics table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Key Performance Indicators"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    # Create table
    rows = len(performance_data['Metric']) + 1  # +1 for header
    cols = len(performance_data)

    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(7)
    height = Inches(3.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(1.5)

    # Add headers
    headers = list(performance_data.keys())
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 78, 120)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Add data rows
    for row_idx in range(len(performance_data['Metric'])):
        for col_idx, key in enumerate(headers):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(performance_data[key][row_idx])
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)


def create_summary_slide(prs):
    """Create a summary slide with key takeaways."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content

    title = slide.shapes.title
    title.text = "Summary & Key Takeaways"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Add bullet points
    points = [
        "Total Annual Revenue: $795K (65% YoY growth)",
        "All performance metrics exceeding targets",
        "Q4 showed strongest performance at $248K",
        "System uptime maintained at 99.95%",
        "Ready to scale for projected 3x traffic growth"
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.space_before = Pt(12)
        p.space_after = Pt(12)


def main():
    """Main function to generate the presentation."""
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Generating presentation...")

    # Create slides
    create_title_slide(prs)
    print("✓ Title slide created")

    create_chart_slide(prs)
    print("✓ Chart slide created")

    create_table_slide(prs)
    print("✓ Table slide created")

    create_summary_slide(prs)
    print("✓ Summary slide created")

    # Save presentation
    output_dir = Path(__file__).parent.parent / 'outputs'
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / 'example_report.pptx'

    prs.save(str(output_path))
    print(f"\n✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
